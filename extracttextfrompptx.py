from pptx import Presentation
import glob

with open("pptxtotext.txt", "w") as f:

 
    for eachfile in glob.glob("*.pptx"):
        powerpoint = Presentation(eachfile)
        
        f.write("{0}".format(str(eachfile)))
        f.write("\n")
        f.write("----------------------\n")
        f.write("\n")
        for slide in powerpoint.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    
                    f.write("{0}".format(str(shape.text)))
                    f.write("\n")
                    